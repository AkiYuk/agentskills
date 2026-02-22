"""PowerPointスライド生成スクリプト

slides.md形式のMarkdownファイルからPowerPoint（.pptx）ファイルを生成する。
18種類のレイアウトに対応し、config.jsonで指定されたデザイン設定を適用する。

使い方:
    python slide_generator_pptx.py --markdown-file output/slides.md --config config.json
    python slide_generator_pptx.py --markdown-file output/slides.md --config config.json --template assets/template.pptx
    python slide_generator_pptx.py --markdown-file output/slides.md --config config.json --output-dir my_output
"""

from __future__ import annotations

import argparse
import json
import os
import sys
from pathlib import Path

from lxml.etree import SubElement, QName

from pptx import Presentation as PptxPresentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION, XL_MARKER_STYLE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from slides_markdown import SlidesMarkdownParser, SlideContent


# ====================================================================
# DesignConfig
# ====================================================================

class DesignConfig:
    """config.jsonからデザイン設定を読み込むクラス"""

    def __init__(self, config_path: str = None):
        """設定を初期化する

        Args:
            config_path: config.jsonファイルのパス。Noneの場合はデフォルト値を使用
        """
        if config_path and os.path.exists(config_path):
            with open(config_path, encoding="utf-8") as f:
                self._config = json.load(f)
        else:
            self._config = {}

    @property
    def font_family(self) -> str:
        return self._config.get("font", {}).get("family", "Meiryo UI")

    @property
    def palette(self) -> dict:
        return self._config.get("palette", {})

    @property
    def primary_color(self) -> str:
        """プライマリカラー (#1E3A5F)"""
        return self.palette.get("primary", "#1E3A5F")

    @property
    def secondary_color(self) -> str:
        """セカンダリカラー (#4A6FA5)"""
        return self.palette.get("secondary", "#4A6FA5")

    @property
    def accent_color(self) -> str:
        """アクセントカラー (#3AA899)"""
        return self.palette.get("accent", "#3AA899")

    @property
    def text_primary(self) -> str:
        """テキスト色 プライマリ (#333333)"""
        return self.palette.get("text", {}).get("primary", "#333333")

    @property
    def text_secondary(self) -> str:
        """テキスト色 セカンダリ (#666666)"""
        return self.palette.get("text", {}).get("secondary", "#666666")

    @property
    def text_light(self) -> str:
        """テキスト色 ライト (#FFFFFF)"""
        return self.palette.get("text", {}).get("light", "#FFFFFF")

    @property
    def bg_primary(self) -> str:
        """背景色 プライマリ (#FFFFFF)"""
        return self.palette.get("background", {}).get("primary", "#FFFFFF")

    @property
    def bg_secondary(self) -> str:
        """背景色 セカンダリ (#F5F5F5)"""
        return self.palette.get("background", {}).get("secondary", "#F5F5F5")

    @property
    def gray_color(self) -> str:
        """グレー (#999999)"""
        return self.palette.get("gray", "#999999")

    @property
    def chart_line_color_3(self) -> str:
        """チャート折れ線3系列目の色 (#EDB120)"""
        return self.palette.get("chart", {}).get("line_color_3", "#EDB120")

    @property
    def bg_dark(self) -> str:
        """背景色 ダーク (#1E3A5F)"""
        return self.palette.get("background", {}).get("dark", "#1E3A5F")

    @property
    def output_dir(self) -> str:
        """出力ディレクトリ"""
        return self._config.get("output", {}).get("dir", "output")

    @property
    def template_path(self) -> str:
        """テンプレートPPTXのパス"""
        return self._config.get("template", {}).get("pptx_path", "")


# ====================================================================
# BaseRenderer
# ====================================================================

class BaseRenderer:
    """全Rendererの基底クラス"""

    # 共通座標定数（Inches）
    MARGIN_LEFT = 0.5
    MARGIN_RIGHT = 0.5
    MARGIN_BOTTOM = 0.5
    CONTENT_LEFT = 0.5
    CONTENT_WIDTH = 12.333  # 13.333 - 0.5 - 0.5

    TITLE_TOP = 0.3
    TITLE_HEIGHT = 0.6
    KEY_MSG_TOP = 0.95
    KEY_MSG_HEIGHT = 0.45
    CONTENT_TOP = 1.55
    CONTENT_HEIGHT = 5.45  # 7.5 - 1.55 - 0.5

    def render(self, slide, slide_data: SlideContent, config: DesignConfig):
        """サブクラスで実装する描画メソッド"""
        raise NotImplementedError

    def _render_slide_title(self, slide, text: str, config: DesignConfig,
                            font_color: str = None):
        """共通: スライドタイトルの描画（24pt）"""
        color = font_color or config.text_primary
        self.add_text_box(
            slide,
            left=Inches(self.CONTENT_LEFT),
            top=Inches(self.TITLE_TOP),
            width=Inches(self.CONTENT_WIDTH),
            height=Inches(self.TITLE_HEIGHT),
            text=text,
            font_size=Pt(24),
            font_color=color,
            bold=True,
            font_family=config.font_family,
        )

    def _render_key_message(self, slide, text: str, config: DesignConfig,
                            font_color: str = None):
        """共通: キーメッセージの描画（16pt）"""
        color = font_color or config.text_secondary
        self.add_text_box(
            slide,
            left=Inches(self.CONTENT_LEFT),
            top=Inches(self.KEY_MSG_TOP),
            width=Inches(self.CONTENT_WIDTH),
            height=Inches(self.KEY_MSG_HEIGHT),
            text=text,
            font_size=Pt(16),
            font_color=color,
            font_family=config.font_family,
        )

    def _render_separator_line(self, slide, config: DesignConfig):
        """共通: タイトル・キーメッセージ下の区切り線"""
        self.add_horizontal_line(
            slide,
            left=Inches(self.CONTENT_LEFT),
            top=Inches(1.45),
            width=Inches(self.CONTENT_WIDTH),
            color="#CCCCCC",
            thickness=1.0,
        )

    def add_text_box(self, slide, left, top, width, height,
                     text: str, font_size=Pt(12),
                     font_color: str = "#333333",
                     bold: bool = False, italic: bool = False,
                     alignment=PP_ALIGN.LEFT,
                     font_family: str = "Meiryo UI",
                     line_spacing: float = 1.5):
        """テキストボックスを追加し、フォント設定を適用"""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = font_size
        p.font.color.rgb = RGBColor.from_string(font_color.lstrip("#"))
        p.font.bold = bold
        p.font.italic = italic
        p.font.name = font_family
        p.alignment = alignment
        p.line_spacing = line_spacing
        return txBox

    def set_background_color(self, slide, hex_color: str):
        """スライドの背景色を単色で設定"""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(hex_color.lstrip("#"))

    def set_background_gradient(self, slide, color1: str, color2: str):
        """グラデーション背景を設定（将来用）"""
        background = slide.background
        fill = background.fill
        fill.gradient()
        fill.gradient_stops[0].color.rgb = RGBColor.from_string(color1.lstrip("#"))
        fill.gradient_stops[1].color.rgb = RGBColor.from_string(color2.lstrip("#"))

    def add_shape(self, slide, shape_type, left, top, width, height,
                  fill_color: str = None, line_color: str = None):
        """図形を追加"""
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor.from_string(fill_color.lstrip("#"))
        if line_color:
            shape.line.color.rgb = RGBColor.from_string(line_color.lstrip("#"))
        else:
            shape.line.fill.background()  # 枠線なし
        return shape

    def add_horizontal_line(self, slide, left, top, width,
                            color: str = "#CCCCCC", thickness: float = 1.0):
        """水平線を描画（細い矩形として実装）"""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, top, width, Pt(thickness)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(color.lstrip("#"))
        shape.line.fill.background()
        return shape

    def add_vertical_line(self, slide, left, top, height,
                          color: str = "#CCCCCC", thickness: float = 1.0):
        """縦線を描画（細い矩形として実装）"""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, top, Pt(thickness), height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(color.lstrip("#"))
        shape.line.fill.background()
        return shape

    def _render_columns(self, slide, columns_data: list, config: DesignConfig,
                        gutter: float, header_font_size=Pt(14),
                        body_font_size=Pt(12)):
        """カラムレイアウト共通の描画処理

        Args:
            slide: スライドオブジェクト
            columns_data: カラムデータのリスト（各要素は heading, body キーを持つdict）
            config: デザイン設定
            gutter: カラム間のガター幅（Inches）
            header_font_size: ヘッダーのフォントサイズ
            body_font_size: 本文のフォントサイズ
        """
        num_cols = len(columns_data)
        col_width = (self.CONTENT_WIDTH - gutter * (num_cols - 1)) / num_cols

        for i, col_data in enumerate(columns_data):
            col_left = Inches(self.CONTENT_LEFT + i * (col_width + gutter))
            # ヘッダー描画
            self.add_text_box(
                slide,
                left=col_left,
                top=Inches(self.CONTENT_TOP),
                width=Inches(col_width),
                height=Inches(0.4),
                text=col_data.get("heading", ""),
                font_size=header_font_size,
                font_color=config.primary_color,
                bold=True,
                font_family=config.font_family,
            )
            # ヘッダー下の装飾線
            self.add_horizontal_line(
                slide,
                left=col_left,
                top=Inches(self.CONTENT_TOP + 0.45),
                width=Inches(col_width),
                color=config.accent_color,
                thickness=2.0,
            )
            # 本文描画
            self.add_text_box(
                slide,
                left=col_left,
                top=Inches(self.CONTENT_TOP + 0.55),
                width=Inches(col_width),
                height=Inches(self.CONTENT_HEIGHT - 0.55),
                text=col_data.get("body", ""),
                font_size=body_font_size,
                font_color=config.text_primary,
                font_family=config.font_family,
            )

    def set_font(self, text_frame, font_family: str = "Meiryo UI",
                 font_size=None, font_color: str = None,
                 bold: bool = None, italic: bool = None,
                 alignment=None, line_spacing: float = None):
        """テキストフレームのフォント設定ヘルパー"""
        for p in text_frame.paragraphs:
            if font_family:
                p.font.name = font_family
            if font_size is not None:
                p.font.size = font_size
            if font_color is not None:
                p.font.color.rgb = RGBColor.from_string(font_color.lstrip("#"))
            if bold is not None:
                p.font.bold = bold
            if italic is not None:
                p.font.italic = italic
            if alignment is not None:
                p.alignment = alignment
            if line_spacing is not None:
                p.line_spacing = line_spacing


# ====================================================================
# TitleRenderer
# ====================================================================

class TitleRenderer(BaseRenderer):
    """表紙レイアウトRenderer

    紺色背景に白文字のメインタイトルとサブタイトルを中央配置する。
    """

    def render(self, slide, slide_data: SlideContent, config: DesignConfig):
        """表紙スライドを描画"""
        # 背景: 紺色
        self.set_background_color(slide, config.bg_dark)

        # メインタイトル: 44pt, bold, 白色, 中央揃え
        self.add_text_box(
            slide,
            left=Inches(1.5),
            top=Inches(2.2),
            width=Inches(10.333),
            height=Inches(1.2),
            text=slide_data.title or "",
            font_size=Pt(44),
            font_color=config.text_light,
            bold=True,
            alignment=PP_ALIGN.CENTER,
            font_family=config.font_family,
        )

        # サブタイトル: 16pt, 白色, 中央揃え
        if slide_data.subtitle:
            self.add_text_box(
                slide,
                left=Inches(1.5),
                top=Inches(3.6),
                width=Inches(10.333),
                height=Inches(0.8),
                text=slide_data.subtitle,
                font_size=Pt(16),
                font_color=config.text_light,
                alignment=PP_ALIGN.CENTER,
                font_family=config.font_family,
            )


# ====================================================================
# SectionRenderer
# ====================================================================

class SectionRenderer(BaseRenderer):
    """セクション区切りレイアウトRenderer

    薄グレー背景にセクションタイトルを中央配置し、アクセント色の水平線で装飾する。
    """

    def render(self, slide, slide_data: SlideContent, config: DesignConfig):
        """セクション区切りスライドを描画"""
        # 背景: 薄グレー
        self.set_background_color(slide, config.bg_secondary)

        # セクションタイトル: 36pt, bold, テキスト色primary, 中央揃え
        self.add_text_box(
            slide,
            left=Inches(1.0),
            top=Inches(2.5),
            width=Inches(11.333),
            height=Inches(1.5),
            text=slide_data.title or "",
            font_size=Pt(36),
            font_color=config.text_primary,
            bold=True,
            alignment=PP_ALIGN.CENTER,
            font_family=config.font_family,
        )

        # アクセント色の水平線装飾
        self.add_horizontal_line(
            slide,
            left=Inches(5.667),
            top=Inches(4.2),
            width=Inches(2.0),
            color=config.accent_color,
            thickness=3.0,
        )


# ====================================================================
# TocRenderer
# ====================================================================

class TocRenderer(BaseRenderer):
    """目次・サマリーレイアウトRenderer

    番号付きの目次項目を表示する。番号部分はアクセント色で装飾する。
    """

    def render(self, slide, slide_data: SlideContent, config: DesignConfig):
        """目次スライドを描画"""
        # スライドタイトルとキーメッセージ
        if slide_data.title:
            self._render_slide_title(slide, slide_data.title, config)
        if slide_data.subtitle:
            self._render_key_message(slide, slide_data.subtitle, config)

        # 区切り線
        self._render_separator_line(slide, config)

        # 目次項目: 番号付きリスト
        items = slide_data.items or []
        if not items:
            return

        # 1つのテキストボックスに複数段落で配置
        txBox = slide.shapes.add_textbox(
            Inches(0.8), Inches(self.CONTENT_TOP),
            Inches(11.733), Inches(self.CONTENT_HEIGHT)
        )
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, item in enumerate(items):
            # 最初の項目は既存のparagraphsを使用、2番目以降はadd_paragraph
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            # 番号Run（accent色, bold）
            num_run = p.add_run()
            num_run.text = f"{i + 1}.  "
            num_run.font.size = Pt(14)
            num_run.font.color.rgb = RGBColor.from_string(config.accent_color.lstrip("#"))
            num_run.font.bold = True
            num_run.font.name = config.font_family

            # テキストRun（text_primary色）
            text_run = p.add_run()
            text_run.text = item.get("title", "")
            text_run.font.size = Pt(14)
            text_run.font.color.rgb = RGBColor.from_string(config.text_primary.lstrip("#"))
            text_run.font.name = config.font_family

            p.line_spacing = 1.5
            p.alignment = PP_ALIGN.LEFT


# ====================================================================
# BulletPointsRenderer
# ====================================================================

class BulletPointsRenderer(BaseRenderer):
    """箇条書きレイアウトRenderer

    ビュレットマーカー付きの箇条書き項目を表示する。マーカーはアクセント色。
    """

    def render(self, slide, slide_data: SlideContent, config: DesignConfig):
        """箇条書きスライドを描画"""
        # スライドタイトルとキーメッセージ
        if slide_data.title:
            self._render_slide_title(slide, slide_data.title, config)
        if slide_data.subtitle:
            self._render_key_message(slide, slide_data.subtitle, config)

        # 区切り線
        self._render_separator_line(slide, config)

        # 箇条書き項目
        items = slide_data.items or []
        if not items:
            return

        # 1つのテキストボックスに複数段落で配置
        txBox = slide.shapes.add_textbox(
            Inches(self.CONTENT_LEFT), Inches(self.CONTENT_TOP),
            Inches(self.CONTENT_WIDTH), Inches(self.CONTENT_HEIGHT)
        )
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            # ビュレットマーカーRun（accent色）
            marker_run = p.add_run()
            marker_run.text = "\u25cf  "  # ●
            marker_run.font.size = Pt(12)
            marker_run.font.color.rgb = RGBColor.from_string(config.accent_color.lstrip("#"))
            marker_run.font.name = config.font_family

            # テキストRun（text_primary色）
            text_run = p.add_run()
            text_run.text = item.get("title", "")
            text_run.font.size = Pt(12)
            text_run.font.color.rgb = RGBColor.from_string(config.text_primary.lstrip("#"))
            text_run.font.name = config.font_family

            p.line_spacing = 1.5
            p.alignment = PP_ALIGN.LEFT


# ====================================================================
# NumberedListRenderer
# ====================================================================

class NumberedListRenderer(BaseRenderer):
    """手順・ステップレイアウトRenderer

    円形ステップ番号（accent色）+ タイトル + 説明で各ステップを表示する。
    """

    def render(self, slide, slide_data: SlideContent, config: DesignConfig):
        """手順・ステップスライドを描画"""
        # スライドタイトルとキーメッセージ
        if slide_data.title:
            self._render_slide_title(slide, slide_data.title, config)
        if slide_data.subtitle:
            self._render_key_message(slide, slide_data.subtitle, config)

        # 区切り線
        self._render_separator_line(slide, config)

        # ステップ項目
        items = slide_data.items or []
        if not items:
            return

        # 項目数に応じてコンテンツ領域を等分
        num_items = len(items)
        step_interval = self.CONTENT_HEIGHT / max(num_items, 1)

        for i, item in enumerate(items):
            step_top = self.CONTENT_TOP + (i * step_interval)

            # ステップ番号の円（accent色）
            circle = self.add_shape(
                slide,
                MSO_SHAPE.OVAL,
                Inches(self.CONTENT_LEFT),
                Inches(step_top),
                Inches(0.45),
                Inches(0.45),
                fill_color=config.accent_color,
            )

            # 円内の番号テキスト（白色, 中央揃え）
            tf = circle.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.text = str(i + 1)
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor.from_string(config.text_light.lstrip("#"))
            p.font.bold = True
            p.font.name = config.font_family
            p.alignment = PP_ALIGN.CENTER

            # ステップタイトル（14pt, bold）
            self.add_text_box(
                slide,
                left=Inches(1.1),
                top=Inches(step_top),
                width=Inches(11.733),
                height=Inches(0.35),
                text=item.get("title", ""),
                font_size=Pt(14),
                font_color=config.text_primary,
                bold=True,
                font_family=config.font_family,
                line_spacing=1.5,
            )

            # ステップ説明（12pt）
            description = item.get("description", "")
            if description:
                self.add_text_box(
                    slide,
                    left=Inches(1.1),
                    top=Inches(step_top + 0.35),
                    width=Inches(11.733),
                    height=Inches(0.4),
                    text=description,
                    font_size=Pt(12),
                    font_color=config.text_secondary,
                    font_family=config.font_family,
                    line_spacing=1.5,
                )


# ====================================================================
# CtaRenderer
# ====================================================================

class CtaRenderer(BaseRenderer):
    """行動喚起レイアウトRenderer

    薄グレー背景にメッセージ・説明・ボタン風テキストを中央配置する。
    """

    def render(self, slide, slide_data: SlideContent, config: DesignConfig):
        """CTAスライドを描画"""
        # 背景: 薄グレー
        self.set_background_color(slide, config.bg_secondary)

        # メッセージ: 24pt, bold, 中央揃え
        self.add_text_box(
            slide,
            left=Inches(1.5),
            top=Inches(1.8),
            width=Inches(10.333),
            height=Inches(0.8),
            text=slide_data.title or "",
            font_size=Pt(24),
            font_color=config.text_primary,
            bold=True,
            alignment=PP_ALIGN.CENTER,
            font_family=config.font_family,
        )

        # 説明: 16pt, テキスト色secondary, 中央揃え
        if slide_data.subtitle:
            self.add_text_box(
                slide,
                left=Inches(1.5),
                top=Inches(2.8),
                width=Inches(10.333),
                height=Inches(0.8),
                text=slide_data.subtitle,
                font_size=Pt(16),
                font_color=config.text_secondary,
                alignment=PP_ALIGN.CENTER,
                font_family=config.font_family,
            )

        # ボタン風テキスト: accent色背景の角丸矩形
        button_text = slide_data.cta_button or ""
        if button_text:
            # ボタン背景（角丸矩形）
            button_shape = self.add_shape(
                slide,
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(4.667),
                Inches(4.2),
                Inches(4.0),
                Inches(0.7),
                fill_color=config.accent_color,
            )
            # 角丸率を設定
            button_shape.adjustments[0] = 0.15

            # ボタンテキスト（図形内テキストフレーム）
            tf = button_shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = button_text
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor.from_string(config.text_light.lstrip("#"))
            p.font.bold = True
            p.font.name = config.font_family
            p.alignment = PP_ALIGN.CENTER


# ====================================================================
# TwoColumnRenderer
# ====================================================================

class TwoColumnRenderer(BaseRenderer):
    """2カラムレイアウトRenderer

    コンテンツ領域を左右均等に2分割し、各カラムにヘッダーと本文を配置する。
    カラム間には縦の区切り線を表示する。
    """

    def render(self, slide, slide_data: SlideContent, config: DesignConfig):
        """2カラムスライドを描画"""
        # 背景: 白色
        self.set_background_color(slide, config.bg_primary)

        # スライドタイトルとキーメッセージ
        if slide_data.title:
            self._render_slide_title(slide, slide_data.title, config)
        if slide_data.subtitle:
            self._render_key_message(slide, slide_data.subtitle, config)

        # カラム描画
        columns = slide_data.columns or []
        if not columns:
            return

        gutter = 0.3
        self._render_columns(slide, columns, config, gutter=gutter,
                             header_font_size=Pt(14), body_font_size=Pt(12))

        # カラム間の縦区切り線
        col_width = (self.CONTENT_WIDTH - gutter) / 2
        line_left = Inches(self.CONTENT_LEFT + col_width + gutter / 2)
        self.add_vertical_line(
            slide,
            left=line_left,
            top=Inches(self.CONTENT_TOP),
            height=Inches(self.CONTENT_HEIGHT),
            color="#CCCCCC",
        )


# ====================================================================
# ThreeColumnRenderer
# ====================================================================

class ThreeColumnRenderer(BaseRenderer):
    """3カラムレイアウトRenderer

    コンテンツ領域を3等分し、各カラムにヘッダーと本文を配置する。
    """

    def render(self, slide, slide_data: SlideContent, config: DesignConfig):
        """3カラムスライドを描画"""
        # 背景: 白色
        self.set_background_color(slide, config.bg_primary)

        # スライドタイトルとキーメッセージ
        if slide_data.title:
            self._render_slide_title(slide, slide_data.title, config)
        if slide_data.subtitle:
            self._render_key_message(slide, slide_data.subtitle, config)

        # カラム描画
        columns = slide_data.columns or []
        if not columns:
            return

        self._render_columns(slide, columns, config, gutter=0.25,
                             header_font_size=Pt(14), body_font_size=Pt(12))


# ====================================================================
# FourColumnRenderer
# ====================================================================

class FourColumnRenderer(BaseRenderer):
    """4カラムレイアウトRenderer

    コンテンツ領域を4等分し、各カラムにヘッダーと本文を配置する。
    フォントサイズをやや小さくしてスペースを確保する。
    """

    def render(self, slide, slide_data: SlideContent, config: DesignConfig):
        """4カラムスライドを描画"""
        # 背景: 白色
        self.set_background_color(slide, config.bg_primary)

        # スライドタイトルとキーメッセージ
        if slide_data.title:
            self._render_slide_title(slide, slide_data.title, config)
        if slide_data.subtitle:
            self._render_key_message(slide, slide_data.subtitle, config)

        # カラム描画（フォントサイズ小さめ）
        columns = slide_data.columns or []
        if not columns:
            return

        self._render_columns(slide, columns, config, gutter=0.2,
                             header_font_size=Pt(12), body_font_size=Pt(11))


# ====================================================================
# MetricsRenderer
# ====================================================================

class MetricsRenderer(BaseRenderer):
    """数値・KPIレイアウトRenderer

    2〜4個の数値カードを横並びで均等配置し、カード間に縦区切り線を表示する。
    """

    def render(self, slide, slide_data: SlideContent, config: DesignConfig):
        """メトリクススライドを描画"""
        # 背景: 白色
        self.set_background_color(slide, config.bg_primary)

        # スライドタイトルとキーメッセージ
        if slide_data.title:
            self._render_slide_title(slide, slide_data.title, config)
        if slide_data.subtitle:
            self._render_key_message(slide, slide_data.subtitle, config)

        metrics = slide_data.items or []
        num = len(metrics)
        if num < 2:
            num = 2
        if num > 4:
            num = 4

        # ガターとカード幅の計算
        gutter = {2: 0.3, 3: 0.25, 4: 0.2}.get(num, 0.25)
        card_width = (self.CONTENT_WIDTH - gutter * (num - 1)) / num

        for i, metric in enumerate(metrics[:4]):
            card_left = Inches(self.CONTENT_LEFT + i * (card_width + gutter))
            # 数値
            self.add_text_box(
                slide,
                left=card_left,
                top=Inches(2.5),
                width=Inches(card_width),
                height=Inches(1.2),
                text=metric.get("value", ""),
                font_size=Pt(36),
                font_color=config.accent_color,
                bold=True,
                alignment=PP_ALIGN.CENTER,
                font_family=config.font_family,
            )
            # ラベル
            self.add_text_box(
                slide,
                left=card_left,
                top=Inches(3.8),
                width=Inches(card_width),
                height=Inches(0.5),
                text=metric.get("label", ""),
                font_size=Pt(12),
                font_color=config.text_secondary,
                alignment=PP_ALIGN.CENTER,
                font_family=config.font_family,
            )
            # カード間区切り線（最後のカード以外）
            if i < num - 1:
                line_left = Inches(
                    self.CONTENT_LEFT + (i + 1) * card_width + i * gutter + gutter / 2
                )
                self.add_vertical_line(
                    slide,
                    left=line_left,
                    top=Inches(2.2),
                    height=Inches(2.5),
                    color="#CCCCCC",
                )


# ====================================================================
# ComparisonTableRenderer
# ====================================================================

class ComparisonTableRenderer(BaseRenderer):
    """比較表レイアウトRenderer

    python-pptxのadd_table APIでテーブルを描画する。
    ヘッダー行は紺色背景・白文字、データ行は交互背景色。
    """

    def render(self, slide, slide_data: SlideContent, config: DesignConfig):
        """比較表スライドを描画"""
        # 背景: 白色
        self.set_background_color(slide, config.bg_primary)

        # スライドタイトルとキーメッセージ
        if slide_data.title:
            self._render_slide_title(slide, slide_data.title, config)
        if slide_data.subtitle:
            self._render_key_message(slide, slide_data.subtitle, config)

        table_data = slide_data.table or []
        if not table_data or len(table_data) < 2:
            return

        headers = table_data[0]
        rows = table_data[1:]
        num_rows = len(rows) + 1  # ヘッダー含む
        num_cols = len(headers)

        # テーブル作成
        table_height = Inches(0.5 + 0.45 * len(rows))
        table_shape = slide.shapes.add_table(
            num_rows, num_cols,
            left=Inches(self.CONTENT_LEFT),
            top=Inches(self.CONTENT_TOP),
            width=Inches(self.CONTENT_WIDTH),
            height=table_height,
        )
        table = table_shape.table

        # ヘッダー行のスタイル設定
        for col_idx, header_text in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header_text
            # 背景色: primary
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor.from_string(
                config.primary_color.lstrip("#"))
            # フォント
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor.from_string("FFFFFF")
                paragraph.font.name = config.font_family
                paragraph.alignment = PP_ALIGN.CENTER
            # セル余白
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.05)
            cell.margin_bottom = Inches(0.05)

        # データ行のスタイル設定
        for row_idx, row_data in enumerate(rows):
            # 0始まりで偶数行が白、奇数行がグレー
            is_even_row = (row_idx % 2 == 1)
            bg_color = config.bg_secondary if is_even_row else config.bg_primary

            for col_idx, cell_text in enumerate(row_data):
                if col_idx >= num_cols:
                    break
                cell = table.cell(row_idx + 1, col_idx)
                # 太字ラベル（**text**）の検出
                is_bold_label = cell_text.startswith("**") and cell_text.endswith("**")
                if is_bold_label:
                    cell_text = cell_text[2:-2]
                cell.text = cell_text
                # 背景色
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor.from_string(bg_color.lstrip("#"))
                # フォント
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(11)
                    paragraph.font.name = config.font_family
                    if is_bold_label:
                        paragraph.font.bold = True
                        paragraph.font.color.rgb = RGBColor.from_string(
                            config.accent_color.lstrip("#"))
                    else:
                        paragraph.font.color.rgb = RGBColor.from_string(
                            config.text_primary.lstrip("#"))
                # セル余白
                cell.margin_left = Inches(0.05)
                cell.margin_right = Inches(0.05)
                cell.margin_top = Inches(0.05)
                cell.margin_bottom = Inches(0.05)


# ====================================================================
# FaqRenderer
# ====================================================================

class FaqRenderer(BaseRenderer):
    """Q&AレイアウトRenderer

    質問と回答のペアを縦並びで表示し、ペア間に水平区切り線を配置する。
    """

    def render(self, slide, slide_data: SlideContent, config: DesignConfig):
        """Q&Aスライドを描画"""
        # 背景: 白色
        self.set_background_color(slide, config.bg_primary)

        # スライドタイトル
        if slide_data.title:
            self._render_slide_title(slide, slide_data.title, config)

        faq_items = slide_data.faq_items or []
        num_pairs = min(len(faq_items), 3)
        if num_pairs == 0:
            return
        pair_height = self.CONTENT_HEIGHT / max(num_pairs, 1)

        for i, item in enumerate(faq_items[:3]):
            pair_top = self.CONTENT_TOP + i * pair_height

            # 質問
            self.add_text_box(
                slide,
                left=Inches(self.CONTENT_LEFT),
                top=Inches(pair_top),
                width=Inches(self.CONTENT_WIDTH),
                height=Inches(0.5),
                text=f"Q.  {item.get('question', '')}",
                font_size=Pt(14),
                font_color=config.primary_color,
                bold=True,
                font_family=config.font_family,
            )

            # 回答
            self.add_text_box(
                slide,
                left=Inches(self.CONTENT_LEFT),
                top=Inches(pair_top + 0.55),
                width=Inches(self.CONTENT_WIDTH),
                height=Inches(pair_height - 0.75),
                text=f"A.  {item.get('answer', '')}",
                font_size=Pt(12),
                font_color=config.text_primary,
                font_family=config.font_family,
            )

            # 区切り線（最終ペア以外）
            if i < num_pairs - 1:
                self.add_horizontal_line(
                    slide,
                    left=Inches(self.CONTENT_LEFT),
                    top=Inches(pair_top + pair_height - 0.15),
                    width=Inches(self.CONTENT_WIDTH),
                    color="#CCCCCC",
                )


# ====================================================================
# QuoteRenderer
# ====================================================================

class QuoteRenderer(BaseRenderer):
    """引用レイアウトRenderer

    左側にaccent色の縦装飾バーを配置し、引用文をitalicで表示する。
    著者名は右寄せで表示する。
    """

    def render(self, slide, slide_data: SlideContent, config: DesignConfig):
        """引用スライドを描画"""
        # 背景: 白色
        self.set_background_color(slide, config.bg_primary)

        # スライドタイトル
        if slide_data.title:
            self._render_slide_title(slide, slide_data.title, config)

        quote_text = slide_data.quote or ""
        author = slide_data.quote_author or ""

        # accent色の装飾バー（縦線）
        self.add_shape(
            slide,
            MSO_SHAPE.RECTANGLE,
            left=Inches(1.0),
            top=Inches(1.8),
            width=Inches(0.06),
            height=Inches(2.5),
            fill_color=config.accent_color,
        )

        # 引用文
        self.add_text_box(
            slide,
            left=Inches(1.4),
            top=Inches(1.8),
            width=Inches(10.933),
            height=Inches(2.5),
            text=quote_text,
            font_size=Pt(16),
            font_color=config.text_primary,
            italic=True,
            font_family=config.font_family,
        )

        # 著者名
        if author:
            self.add_text_box(
                slide,
                left=Inches(1.4),
                top=Inches(4.5),
                width=Inches(10.933),
                height=Inches(0.5),
                text=f"-- {author}",
                font_size=Pt(12),
                font_color=config.text_secondary,
                alignment=PP_ALIGN.RIGHT,
                font_family=config.font_family,
            )


# ====================================================================
# ImageWithTextRenderer
# ====================================================================

class ImageWithTextRenderer(BaseRenderer):
    """画像+テキストレイアウトRenderer

    左半分に画像、右半分にテキスト説明を配置する。
    画像はアスペクト比を維持してフィットさせる。
    """

    def render(self, slide, slide_data: SlideContent, config: DesignConfig):
        """画像+テキストスライドを描画"""
        # 背景: 白色
        self.set_background_color(slide, config.bg_primary)

        # スライドタイトルとキーメッセージ
        if slide_data.title:
            self._render_slide_title(slide, slide_data.title, config)
        if slide_data.subtitle:
            self._render_key_message(slide, slide_data.subtitle, config)

        image_path = slide_data.image_path or ""
        description = slide_data.body or ""

        # 画像エリアとテキストエリアのサイズ
        area_gutter = 0.4
        area_width = (self.CONTENT_WIDTH - area_gutter) / 2  # 5.967in

        img_left = Inches(self.CONTENT_LEFT)
        img_top = Inches(self.CONTENT_TOP)
        img_area_width = Inches(area_width)
        img_area_height = Inches(self.CONTENT_HEIGHT)

        if image_path and os.path.exists(image_path):
            # アスペクト比を維持してフィット
            try:
                from PIL import Image
                with Image.open(image_path) as img:
                    img_w, img_h = img.size
                aspect = img_w / img_h
                target_aspect = area_width / self.CONTENT_HEIGHT

                if aspect > target_aspect:
                    # 横長: 幅に合わせ、高さを調整
                    actual_width = img_area_width
                    actual_height = Inches(area_width / aspect)
                    actual_top = img_top + (img_area_height - actual_height) // 2
                    actual_left = img_left
                else:
                    # 縦長: 高さに合わせ、幅を調整
                    actual_height = img_area_height
                    actual_width = Inches(self.CONTENT_HEIGHT * aspect)
                    actual_left = img_left + (img_area_width - actual_width) // 2
                    actual_top = img_top

                slide.shapes.add_picture(
                    image_path, actual_left, actual_top,
                    actual_width, actual_height
                )
            except ImportError:
                # PILがない場合は固定サイズで配置
                slide.shapes.add_picture(
                    image_path, img_left, img_top,
                    img_area_width, img_area_height
                )
        else:
            # 画像なし: placeholder.pngを使用
            placeholder_path = os.path.join(
                os.path.dirname(__file__), "..", "assets", "placeholder.png"
            )
            if os.path.exists(placeholder_path):
                slide.shapes.add_picture(
                    placeholder_path, img_left, img_top,
                    img_area_width, img_area_height
                )
            else:
                # プレースホルダー画像もない場合はグレー矩形で代替
                self.add_shape(
                    slide, MSO_SHAPE.RECTANGLE,
                    img_left, img_top, img_area_width, img_area_height,
                    fill_color="#E0E0E0",
                )

        # テキスト説明（右半分）
        text_left = Inches(self.CONTENT_LEFT + area_width + area_gutter)
        self.add_text_box(
            slide,
            left=text_left,
            top=Inches(self.CONTENT_TOP),
            width=Inches(area_width),
            height=Inches(self.CONTENT_HEIGHT),
            text=description,
            font_size=Pt(12),
            font_color=config.text_primary,
            font_family=config.font_family,
        )


# ====================================================================
# ShapeHelpers
# ====================================================================

# OOXML名前空間
_NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_NS_C = "http://schemas.openxmlformats.org/drawingml/2006/chart"


class ShapeHelpers:
    """python-pptxのAPIでは表現できないデザインをOOXML直接操作で実現するヘルパー"""

    @staticmethod
    def set_bar_shape_rounded(chart) -> None:
        """バーの形状を角丸に設定する（OOXML操作）

        各系列の<c:spPr>に<a:prstGeom prst="roundRect"/>を追加する。
        """
        chart_elem = chart.element
        # barChart または bar3DChart 内の ser 要素を取得
        for tag in ("barChart", "bar3DChart"):
            for ser in chart_elem.iter(QName(_NS_C, "ser")):
                # spPr要素を取得または作成
                sp_pr = ser.find(QName(_NS_C, "spPr"))
                if sp_pr is None:
                    sp_pr = ser.find(QName(_NS_A, "spPr"))
                if sp_pr is None:
                    sp_pr = SubElement(ser, QName(_NS_A, "spPr"))
                # 既存のprstGeomを除去
                for existing in sp_pr.findall(QName(_NS_A, "prstGeom")):
                    sp_pr.remove(existing)
                # 角丸矩形を設定
                prst_geom = SubElement(sp_pr, QName(_NS_A, "prstGeom"))
                prst_geom.set("prst", "roundRect")

    @staticmethod
    def set_gridline_dashed(axis) -> None:
        """軸のグリッド線を破線スタイル・#E5E5E5に設定する（OOXML操作）"""
        # メジャーグリッド線を有効化
        axis.has_major_gridlines = True

        # 軸XML要素からmajorGridlinesを取得
        axis_elem = axis._element
        gridlines_elem = axis_elem.find(QName(_NS_C, "majorGridlines"))
        if gridlines_elem is None:
            return

        # spPrを取得または作成
        sp_pr = gridlines_elem.find(QName(_NS_C, "spPr"))
        if sp_pr is None:
            sp_pr = gridlines_elem.find(QName(_NS_A, "spPr"))
        if sp_pr is None:
            sp_pr = SubElement(gridlines_elem, QName(_NS_A, "spPr"))

        # 既存のlnを除去
        for existing_ln in sp_pr.findall(QName(_NS_A, "ln")):
            sp_pr.remove(existing_ln)

        # ln要素を作成
        ln = SubElement(sp_pr, QName(_NS_A, "ln"))

        # 色を設定（#E5E5E5）
        solid_fill = SubElement(ln, QName(_NS_A, "solidFill"))
        srgb_clr = SubElement(solid_fill, QName(_NS_A, "srgbClr"))
        srgb_clr.set("val", "E5E5E5")

        # 破線スタイルを設定
        prst_dash = SubElement(ln, QName(_NS_A, "prstDash"))
        prst_dash.set("val", "dash")

    @staticmethod
    def set_line_marker(series, color_hex: str, marker_size: int = 8) -> None:
        """折れ線グラフのマーカーを設定する（APIとOOXML併用）

        円形マーカー、白塗り、枠線は系列色と同色・1.5pt。
        """
        # マーカースタイルとサイズ（API）
        series.marker.style = XL_MARKER_STYLE.CIRCLE
        series.marker.size = marker_size

        # 白塗り
        series.marker.format.fill.solid()
        series.marker.format.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        # 枠線色
        color_hex_clean = color_hex.lstrip("#")
        series.marker.format.line.color.rgb = RGBColor.from_string(color_hex_clean)
        series.marker.format.line.width = Pt(1.5)

    @staticmethod
    def set_bar_color(chart, series_index: int, point_index: int, color_hex: str) -> None:
        """個別のバーポイントの色を設定する"""
        color_hex_clean = color_hex.lstrip("#")
        point = chart.series[series_index].points[point_index]
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = RGBColor.from_string(color_hex_clean)

    @staticmethod
    def remove_chart_border(chart) -> None:
        """チャートのプロットエリア枠線を非表示にする（OOXML操作）"""
        chart_elem = chart.element
        plot_area = chart_elem.find(QName(_NS_C, "plotArea"))
        if plot_area is None:
            return
        # spPrを取得または作成
        sp_pr = plot_area.find(QName(_NS_C, "spPr"))
        if sp_pr is None:
            sp_pr = plot_area.find(QName(_NS_A, "spPr"))
        if sp_pr is None:
            sp_pr = SubElement(plot_area, QName(_NS_A, "spPr"))
        # 既存のlnを除去
        for existing_ln in sp_pr.findall(QName(_NS_A, "ln")):
            sp_pr.remove(existing_ln)
        # noFillの枠線を設定
        ln = SubElement(sp_pr, QName(_NS_A, "ln"))
        SubElement(ln, QName(_NS_A, "noFill"))


# ====================================================================
# ChartRenderer
# ====================================================================

class ChartRenderer(BaseRenderer):
    """チャートレイアウト（横棒・縦棒・折れ線・100%積み上げ棒）のRenderer"""

    # チャート配置領域
    CHART_LEFT = 0.8
    CHART_TOP = 2.0
    CHART_WIDTH = 11.5
    CHART_HEIGHT = 5.0

    def render(self, slide, slide_data: SlideContent, config: DesignConfig):
        """chart_typeに応じてサブメソッドにディスパッチする"""
        # タイトル・キーメッセージはBaseRendererの共通処理
        if slide_data.title:
            self._render_slide_title(slide, slide_data.title, config)
        if slide_data.subtitle:
            self._render_key_message(slide, slide_data.subtitle, config)

        chart_data = slide_data.chart
        if chart_data is None:
            return

        chart_type = chart_data.chart_type
        if chart_type == "horizontal_bar":
            self._render_horizontal_bar(slide, slide_data, config)
        elif chart_type == "bar":
            self._render_bar(slide, slide_data, config)
        elif chart_type == "line":
            self._render_line(slide, slide_data, config)
        elif chart_type == "pie":
            self._render_pie(slide, slide_data, config)

    def _render_horizontal_bar(self, slide, slide_data: SlideContent,
                                config: DesignConfig) -> None:
        """横棒グラフを描画する"""
        chart_data = slide_data.chart

        # CategoryChartDataを作成
        cd = CategoryChartData()
        cd.categories = chart_data.labels
        # 単一系列（最初の系列のみ使用）
        series_info = chart_data.series[0] if chart_data.series else {"name": "", "values": []}
        cd.add_series(series_info["name"], series_info["values"])

        # チャートを追加
        chart_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED,
            Inches(self.CHART_LEFT), Inches(self.CHART_TOP),
            Inches(self.CHART_WIDTH), Inches(self.CHART_HEIGHT),
            cd
        )
        chart = chart_frame.chart

        # 角丸バー
        ShapeHelpers.set_bar_shape_rounded(chart)

        # バーの色設定（通常=グレー、太字ラベル=アクセント色）
        num_points = len(chart_data.labels)
        for i in range(num_points):
            if i in chart_data.highlight_indices:
                ShapeHelpers.set_bar_color(chart, 0, i, config.accent_color)
            else:
                ShapeHelpers.set_bar_color(chart, 0, i, config.gray_color)

        # 値軸（横軸）を非表示
        chart.value_axis.visible = False

        # カテゴリ軸（縦軸）のフォント設定
        cat_axis = chart.category_axis
        cat_axis.tick_labels.font.size = Pt(10)
        cat_axis.tick_labels.font.color.rgb = RGBColor.from_string(
            config.text_primary.lstrip("#"))
        cat_axis.tick_labels.font.name = config.font_family

        # 破線グリッド線
        ShapeHelpers.set_gridline_dashed(chart.value_axis)

        # データラベル
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.font.size = Pt(9)
        data_labels.font.name = config.font_family
        data_labels.font.color.rgb = RGBColor.from_string(
            config.text_primary.lstrip("#"))
        data_labels.number_format = '0'
        data_labels.show_value = True

        # 凡例非表示
        chart.has_legend = False

        # プロットエリア枠線非表示
        ShapeHelpers.remove_chart_border(chart)

    def _render_bar(self, slide, slide_data: SlideContent,
                     config: DesignConfig) -> None:
        """縦棒グラフを描画する"""
        chart_data = slide_data.chart

        # CategoryChartDataを作成
        cd = CategoryChartData()
        cd.categories = chart_data.labels
        series_info = chart_data.series[0] if chart_data.series else {"name": "", "values": []}
        cd.add_series(series_info["name"], series_info["values"])

        # チャートを追加
        chart_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(self.CHART_LEFT), Inches(self.CHART_TOP),
            Inches(self.CHART_WIDTH), Inches(self.CHART_HEIGHT),
            cd
        )
        chart = chart_frame.chart

        # 角丸バー
        ShapeHelpers.set_bar_shape_rounded(chart)

        # バーの色設定（通常=グレー、太字ラベル=アクセント色）
        num_points = len(chart_data.labels)
        for i in range(num_points):
            if i in chart_data.highlight_indices:
                ShapeHelpers.set_bar_color(chart, 0, i, config.accent_color)
            else:
                ShapeHelpers.set_bar_color(chart, 0, i, config.gray_color)

        # カテゴリ軸（横軸）のフォント設定
        cat_axis = chart.category_axis
        cat_axis.tick_labels.font.size = Pt(10)
        cat_axis.tick_labels.font.color.rgb = RGBColor.from_string(
            config.text_primary.lstrip("#"))
        cat_axis.tick_labels.font.name = config.font_family

        # 値軸（縦軸）のフォント設定
        val_axis = chart.value_axis
        val_axis.tick_labels.font.size = Pt(9)
        val_axis.tick_labels.font.color.rgb = RGBColor.from_string(
            config.text_secondary.lstrip("#"))
        val_axis.tick_labels.font.name = config.font_family

        # 破線グリッド線（値軸のメジャーグリッド線）
        ShapeHelpers.set_gridline_dashed(chart.value_axis)

        # データラベル（バー上部に値を表示）
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.font.size = Pt(9)
        data_labels.font.name = config.font_family
        data_labels.font.color.rgb = RGBColor.from_string(
            config.text_primary.lstrip("#"))
        data_labels.number_format = '0'
        data_labels.show_value = True
        data_labels.label_position = XL_LABEL_POSITION.OUTSIDE_END

        # 凡例非表示
        chart.has_legend = False

        # プロットエリア枠線非表示
        ShapeHelpers.remove_chart_border(chart)

    def _render_line(self, slide, slide_data: SlideContent,
                      config: DesignConfig) -> None:
        """折れ線グラフを描画する"""
        chart_data = slide_data.chart

        # CategoryChartDataを作成（複数系列対応）
        cd = CategoryChartData()
        cd.categories = chart_data.labels
        for s in chart_data.series:
            cd.add_series(s["name"], s["values"])

        # チャートを追加
        chart_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.LINE_MARKERS,
            Inches(self.CHART_LEFT), Inches(self.CHART_TOP),
            Inches(self.CHART_WIDTH), Inches(self.CHART_HEIGHT),
            cd
        )
        chart = chart_frame.chart

        # 系列ごとの色設定
        line_colors = [
            config.primary_color,      # 系列1: primary (#1E3A5F)
            config.accent_color,       # 系列2: accent (#3AA899)
            config.chart_line_color_3,  # 系列3: line_color_3 (#EDB120)
            config.secondary_color,    # 系列4: secondary (#4A6FA5)
        ]

        for idx, series in enumerate(chart.series):
            # 色の決定
            if idx < len(line_colors):
                color_hex = line_colors[idx]
            else:
                color_hex = config.secondary_color
            color_hex_clean = color_hex.lstrip("#")

            # 線の太さと色
            series.format.line.width = Pt(2.5)
            series.format.line.color.rgb = RGBColor.from_string(color_hex_clean)

            # マーカー設定
            ShapeHelpers.set_line_marker(series, color_hex, marker_size=8)

        # カテゴリ軸のフォント設定
        cat_axis = chart.category_axis
        cat_axis.tick_labels.font.size = Pt(10)
        cat_axis.tick_labels.font.color.rgb = RGBColor.from_string(
            config.text_primary.lstrip("#"))
        cat_axis.tick_labels.font.name = config.font_family

        # 値軸のフォント設定
        val_axis = chart.value_axis
        val_axis.tick_labels.font.size = Pt(9)
        val_axis.tick_labels.font.color.rgb = RGBColor.from_string(
            config.text_secondary.lstrip("#"))
        val_axis.tick_labels.font.name = config.font_family

        # 破線グリッド線（値軸のメジャーグリッド線）
        ShapeHelpers.set_gridline_dashed(chart.value_axis)

        # 凡例（系列が2つ以上の場合にチャート下部に表示）
        if len(chart_data.series) >= 2:
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
            chart.legend.font.size = Pt(9)
            chart.legend.font.name = config.font_family
        else:
            chart.has_legend = False

        # プロットエリア枠線非表示
        ShapeHelpers.remove_chart_border(chart)

    def _render_pie(self, slide, slide_data: SlideContent,
                     config: DesignConfig) -> None:
        """100%積み上げ棒グラフ（pieタイプ）を描画する

        各ラベルを系列として扱い、1つの固定カテゴリに対して積み上げる。
        """
        chart_data = slide_data.chart

        # 色パレット（primary, accent, secondary, gray の順）
        base_colors = [
            config.primary_color,
            config.accent_color,
            config.secondary_color,
            config.gray_color,
        ]

        labels = chart_data.labels
        values = chart_data.series[0]["values"] if chart_data.series else []
        highlight_set = set(chart_data.highlight_indices)

        # 色割り当て: 太字ラベルにはアクセント色を優先
        series_colors = []
        non_highlight_color_idx = 0
        # 太字ラベル以外に使う色リスト（accentを除外）
        non_highlight_colors = [c for c in base_colors if c != config.accent_color]

        for i in range(len(labels)):
            if i in highlight_set:
                series_colors.append(config.accent_color)
            else:
                color = non_highlight_colors[non_highlight_color_idx % len(non_highlight_colors)]
                series_colors.append(color)
                non_highlight_color_idx += 1

        # CategoryChartDataを作成
        # カテゴリ=固定1項目、系列=各データ項目
        cd = CategoryChartData()
        cd.categories = [""]  # 固定1カテゴリ
        for i, label in enumerate(labels):
            val = values[i] if i < len(values) else 0
            cd.add_series(label, [val])

        # チャートを追加
        chart_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_STACKED_100,
            Inches(self.CHART_LEFT), Inches(self.CHART_TOP),
            Inches(self.CHART_WIDTH), Inches(self.CHART_HEIGHT),
            cd
        )
        chart = chart_frame.chart

        # 各系列の色を設定
        for idx, series in enumerate(chart.series):
            color_hex = series_colors[idx] if idx < len(series_colors) else config.gray_color
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = RGBColor.from_string(
                color_hex.lstrip("#"))

        # ギャップ幅の縮小
        chart.plots[0].gap_width = 50

        # カテゴリ軸・値軸を非表示
        chart.category_axis.visible = False
        chart.value_axis.visible = False

        # データラベル: 各系列のポイントに%値を表示（白色、10pt）
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.font.size = Pt(10)
        data_labels.font.name = config.font_family
        data_labels.font.color.rgb = RGBColor.from_string(
            config.text_light.lstrip("#"))
        data_labels.number_format = '0"%"'
        data_labels.show_value = True
        data_labels.show_category_name = False
        data_labels.show_series_name = False

        # 凡例（チャート下部に表示）
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(9)
        chart.legend.font.name = config.font_family

        # プロットエリア枠線非表示
        ShapeHelpers.remove_chart_border(chart)


# ====================================================================
# TimelineRenderer
# ====================================================================

class TimelineRenderer(BaseRenderer):
    """タイムライン・ロードマップレイアウトRenderer

    横方向のタイムライン線にマーカーを配置し、時期ラベルと内容を表示する。
    """

    def render(self, slide, slide_data: SlideContent, config: DesignConfig):
        """タイムラインスライドを描画"""
        self.set_background_color(slide, config.bg_primary)

        # タイトル・キーメッセージ
        if slide_data.title:
            self._render_slide_title(slide, slide_data.title, config)
        if slide_data.subtitle:
            self._render_key_message(slide, slide_data.subtitle, config)

        # 区切り線
        self._render_separator_line(slide, config)

        items = slide_data.items or []
        if not items:
            return

        num = min(len(items), 5)
        items = items[:num]

        # タイムラインの中心Y座標
        line_y = 3.8
        # タイムライン水平線
        line_margin = 1.0
        line_width = self.CONTENT_WIDTH - line_margin * 2 + 1.0
        self.add_horizontal_line(
            slide,
            left=Inches(self.CONTENT_LEFT + line_margin - 0.5),
            top=Inches(line_y),
            width=Inches(line_width),
            color=config.gray_color,
            thickness=2.0,
        )

        # 各マーカーの間隔を計算
        usable_width = self.CONTENT_WIDTH - line_margin * 2
        spacing = usable_width / (num - 1) if num > 1 else 0

        for i, item in enumerate(items):
            # マーカーのX座標（中心）
            marker_cx = self.CONTENT_LEFT + line_margin + i * spacing
            marker_size = 0.3

            # 円形マーカー（アクセント色）
            self.add_shape(
                slide, MSO_SHAPE.OVAL,
                Inches(marker_cx - marker_size / 2),
                Inches(line_y - marker_size / 2 + 0.02),
                Inches(marker_size), Inches(marker_size),
                fill_color=config.accent_color,
            )

            # 時期ラベル（マーカーの上）
            label_width = 1.8
            self.add_text_box(
                slide,
                left=Inches(marker_cx - label_width / 2),
                top=Inches(line_y - 1.0),
                width=Inches(label_width),
                height=Inches(0.5),
                text=item.get("label", ""),
                font_size=Pt(11),
                font_color=config.primary_color,
                bold=True,
                alignment=PP_ALIGN.CENTER,
                font_family=config.font_family,
            )

            # 内容テキスト（マーカーの下）
            body_width = 2.0
            self.add_text_box(
                slide,
                left=Inches(marker_cx - body_width / 2),
                top=Inches(line_y + 0.5),
                width=Inches(body_width),
                height=Inches(1.5),
                text=item.get("body", ""),
                font_size=Pt(10),
                font_color=config.text_primary,
                alignment=PP_ALIGN.CENTER,
                font_family=config.font_family,
                line_spacing=1.3,
            )


# ====================================================================
# ThankYouRenderer
# ====================================================================

class ThankYouRenderer(BaseRenderer):
    """終了スライドレイアウトRenderer

    紺色背景に白文字でメッセージを中央配置し、アクセント色の装飾線を添える。
    """

    def render(self, slide, slide_data: SlideContent, config: DesignConfig):
        """終了スライドを描画"""
        # 背景: 紺色（titleと同系統）
        self.set_background_color(slide, config.bg_dark)

        # メインメッセージ: 36pt, bold, 白色, 中央揃え
        self.add_text_box(
            slide,
            left=Inches(1.5),
            top=Inches(2.2),
            width=Inches(10.333),
            height=Inches(1.2),
            text=slide_data.title or "",
            font_size=Pt(36),
            font_color=config.text_light,
            bold=True,
            alignment=PP_ALIGN.CENTER,
            font_family=config.font_family,
        )

        # アクセント色の水平装飾線
        self.add_horizontal_line(
            slide,
            left=Inches(5.667),
            top=Inches(3.6),
            width=Inches(2.0),
            color=config.accent_color,
            thickness=3.0,
        )

        # サブテキスト: 16pt, 白色, 中央揃え（連絡先等）
        if slide_data.subtitle:
            self.add_text_box(
                slide,
                left=Inches(1.5),
                top=Inches(4.0),
                width=Inches(10.333),
                height=Inches(0.8),
                text=slide_data.subtitle,
                font_size=Pt(16),
                font_color=config.text_light,
                alignment=PP_ALIGN.CENTER,
                font_family=config.font_family,
            )


# ====================================================================
# MatrixRenderer
# ====================================================================

class MatrixRenderer(BaseRenderer):
    """2x2マトリクスレイアウトRenderer

    SWOT分析やポジショニングマップ等、4象限で情報を整理する。
    十字の区切り線で4つのセルに分割し、各セルにヘッダと本文を配置する。
    """

    def render(self, slide, slide_data: SlideContent, config: DesignConfig):
        """2x2マトリクススライドを描画"""
        self.set_background_color(slide, config.bg_primary)

        # タイトル・キーメッセージ
        if slide_data.title:
            self._render_slide_title(slide, slide_data.title, config)
        if slide_data.subtitle:
            self._render_key_message(slide, slide_data.subtitle, config)

        # 区切り線
        self._render_separator_line(slide, config)

        columns = slide_data.columns or []
        # 4象限に満たない場合は空で埋める
        while len(columns) < 4:
            columns.append({"heading": "", "body": ""})

        # グリッド座標
        grid_top = self.CONTENT_TOP + 0.1
        grid_height = self.CONTENT_HEIGHT - 0.1
        cell_width = (self.CONTENT_WIDTH - 0.3) / 2  # 中央のガター分を引く
        cell_height = (grid_height - 0.3) / 2  # 中央のガター分を引く

        # 十字の区切り線
        center_x = self.CONTENT_LEFT + self.CONTENT_WIDTH / 2
        center_y = grid_top + grid_height / 2

        # 縦線
        self.add_vertical_line(
            slide,
            left=Inches(center_x),
            top=Inches(grid_top),
            height=Inches(grid_height),
            color="#CCCCCC",
        )
        # 横線
        self.add_horizontal_line(
            slide,
            left=Inches(self.CONTENT_LEFT),
            top=Inches(center_y),
            width=Inches(self.CONTENT_WIDTH),
            color="#CCCCCC",
        )

        # 4象限の配置: 左上→右上→左下→右下
        positions = [
            (self.CONTENT_LEFT, grid_top),                            # 左上
            (self.CONTENT_LEFT + cell_width + 0.3, grid_top),         # 右上
            (self.CONTENT_LEFT, grid_top + cell_height + 0.3),        # 左下
            (self.CONTENT_LEFT + cell_width + 0.3, grid_top + cell_height + 0.3),  # 右下
        ]

        for i, (col_data, (px, py)) in enumerate(zip(columns[:4], positions)):
            # ヘッダー
            self.add_text_box(
                slide,
                left=Inches(px),
                top=Inches(py),
                width=Inches(cell_width),
                height=Inches(0.4),
                text=col_data.get("heading", ""),
                font_size=Pt(14),
                font_color=config.primary_color,
                bold=True,
                font_family=config.font_family,
            )
            # ヘッダー下の装飾線
            self.add_horizontal_line(
                slide,
                left=Inches(px),
                top=Inches(py + 0.45),
                width=Inches(cell_width),
                color=config.accent_color,
                thickness=2.0,
            )
            # 本文
            self.add_text_box(
                slide,
                left=Inches(px),
                top=Inches(py + 0.55),
                width=Inches(cell_width),
                height=Inches(cell_height - 0.6),
                text=col_data.get("body", ""),
                font_size=Pt(11),
                font_color=config.text_primary,
                font_family=config.font_family,
            )


# ====================================================================
# RENDERER_MAP
# ====================================================================

RENDERER_MAP = {
    "title": TitleRenderer,
    "section": SectionRenderer,
    "toc": TocRenderer,
    "bullet_points": BulletPointsRenderer,
    "numbered_list": NumberedListRenderer,
    "cta": CtaRenderer,
    "two_column": TwoColumnRenderer,
    "three_column": ThreeColumnRenderer,
    "four_column": FourColumnRenderer,
    "metrics": MetricsRenderer,
    "comparison_table": ComparisonTableRenderer,
    "faq": FaqRenderer,
    "quote": QuoteRenderer,
    "image_with_text": ImageWithTextRenderer,
    "chart": ChartRenderer,
    "timeline": TimelineRenderer,
    "thank_you": ThankYouRenderer,
    "matrix": MatrixRenderer,
}


# ====================================================================
# テンプレート処理・プレゼンテーション生成
# ====================================================================

def create_presentation(config: DesignConfig, template_path: str = None) -> PptxPresentation:
    """テンプレートの有無に応じてPresentationオブジェクトを生成する

    Args:
        config: デザイン設定
        template_path: テンプレートPPTXファイルのパス

    Returns:
        python-pptxのPresentationオブジェクト
    """
    if template_path and os.path.exists(template_path):
        prs = PptxPresentation(template_path)
        # 既存スライドを全削除（逆順で削除）
        for i in range(len(prs.slides) - 1, -1, -1):
            rId = prs.slides._sldIdLst[i].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[i]
    else:
        prs = PptxPresentation()

    # スライドサイズを16:9に設定
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def add_blank_slide(prs: PptxPresentation):
    """空白スライドを追加し、プレースホルダーを除去する

    Args:
        prs: python-pptxのPresentationオブジェクト

    Returns:
        追加されたスライドオブジェクト
    """
    slide_layout = prs.slide_layouts[6]  # 空白レイアウト（インデックス6）
    slide = prs.slides.add_slide(slide_layout)
    # プレースホルダーを除去
    for ph in list(slide.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)
    return slide


def resolve_template_path(args_template: str, config: DesignConfig) -> str:
    """テンプレートPPTXのパスを優先順位に従って解決する

    優先順位: CLI --template > config.json template.pptx_path > なし

    Args:
        args_template: CLIで指定されたテンプレートパス
        config: デザイン設定

    Returns:
        テンプレートファイルのパス（なしの場合はNone）
    """
    if args_template:
        return args_template
    config_template = config.template_path
    if config_template and os.path.exists(config_template):
        return config_template
    return None


def generate_presentation(markdown_path: str, config: DesignConfig,
                          template_path: str = None, output_dir: str = "output",
                          title: str = None, layout_filter: str = None) -> str:
    """slides.mdからPPTXを生成するメイン関数

    Args:
        markdown_path: slides.mdファイルのパス
        config: デザイン設定
        template_path: テンプレートPPTXのパス
        output_dir: 出力ディレクトリ
        title: 出力ファイル名のプレフィックス
        layout_filter: 特定レイアウトのみ生成（デバッグ用）

    Returns:
        生成されたPPTXファイルのパス
    """
    # 1. slides.mdをパース
    try:
        with open(markdown_path, encoding="utf-8") as f:
            markdown_text = f.read()
    except FileNotFoundError:
        print(f"エラー: ファイルが見つかりません: {markdown_path}", file=sys.stderr)
        sys.exit(1)
    except OSError as e:
        print(f"エラー: ファイルの読み込みに失敗しました: {e}", file=sys.stderr)
        sys.exit(1)

    parser = SlidesMarkdownParser()
    presentation_data = parser.parse(markdown_text)

    # スライドが0件の場合はエラー
    if not presentation_data.slides:
        print("エラー: スライドが見つかりません", file=sys.stderr)
        sys.exit(1)

    # 2. Presentationオブジェクト生成
    prs = create_presentation(config, template_path)

    # 3. タイトルの決定
    output_title = title or presentation_data.title or "presentation"

    # 4. 各スライドを描画
    for slide_data in presentation_data.slides:
        layout = slide_data.layout or "bullet_points"
        if layout_filter and layout != layout_filter:
            continue
        renderer_cls = RENDERER_MAP.get(layout)
        if renderer_cls is None:
            print(f"警告: 未対応レイアウト '{layout}' をスキップします")
            continue
        slide = add_blank_slide(prs)
        renderer_cls().render(slide, slide_data, config)

    # 5. 保存
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, f"{output_title}.pptx")
    prs.save(output_path)
    return output_path


# ====================================================================
# CLI
# ====================================================================

def main():
    """CLIエントリポイント"""
    parser = argparse.ArgumentParser(
        description="slides.mdからPowerPointファイルを生成する"
    )
    parser.add_argument(
        "--markdown-file",
        type=str,
        required=True,
        help="slides.mdファイルのパス（必須）",
    )
    parser.add_argument(
        "--config",
        type=str,
        default=None,
        help="config.jsonファイルのパス（デフォルト: skills/slide-generator/config.json）",
    )
    parser.add_argument(
        "--title",
        type=str,
        default=None,
        help="出力ファイル名のプレフィックス（オプション）",
    )
    parser.add_argument(
        "--output-dir",
        type=str,
        default=None,
        help="出力ディレクトリ（デフォルト: config.jsonのoutput.dir）",
    )
    parser.add_argument(
        "--template",
        type=str,
        default=None,
        help="テンプレートPPTXのパス（オプション）",
    )
    parser.add_argument(
        "--layout",
        type=str,
        default=None,
        help="特定レイアウトのみを生成（デバッグ用）",
    )

    args = parser.parse_args()

    # config.jsonのパスを解決
    config_path = args.config
    if config_path is None:
        # デフォルト: スクリプトの親ディレクトリのconfig.json
        script_dir = Path(__file__).resolve().parent
        default_config = script_dir.parent / "config.json"
        if default_config.exists():
            config_path = str(default_config)
    elif not os.path.exists(config_path):
        # 明示的に指定されたconfig.jsonが存在しない場合はエラー
        print(f"エラー: 設定ファイルが見つかりません: {config_path}", file=sys.stderr)
        sys.exit(1)

    # DesignConfigの読み込み
    config = DesignConfig(config_path)

    # 出力ディレクトリの決定
    output_dir = args.output_dir or config.output_dir

    # テンプレートパスの解決
    template_path = resolve_template_path(args.template, config)

    # PPTX生成
    output_path = generate_presentation(
        markdown_path=args.markdown_file,
        config=config,
        template_path=template_path,
        output_dir=output_dir,
        title=args.title,
        layout_filter=args.layout,
    )

    print(f"生成完了: {output_path}")


if __name__ == "__main__":
    main()
